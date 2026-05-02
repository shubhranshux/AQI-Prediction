"""Generate AQI Prediction Seminar PPT."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

BASE = os.path.dirname(__file__)
OUT = os.path.join(BASE, 'AQI_Prediction_Seminar_v2.pptx')

# Colors
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0x00, 0xD2, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
CARD_BG = RGBColor(0x16, 0x21, 0x3E)
GREEN = RGBColor(0x00, 0xE6, 0x76)
ORANGE = RGBColor(0xFF, 0x99, 0x33)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tf


def add_bullet_slide(slide, left, top, width, height, items, size=16, color=WHITE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
    return tf


def add_accent_line(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def add_section_header(slide, title):
    add_bg(slide)
    add_accent_line(slide, Inches(0.8), Inches(0.9), Inches(2))
    add_text(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(1), title, size=36, color=ACCENT, bold=True)


# ── Slide 1: Title ──────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
         "AQI PREDICTION SYSTEM", size=44, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
         "Machine Learning-Based Air Quality Index Prediction for Odisha",
         size=22, color=WHITE, align=PP_ALIGN.CENTER)
add_accent_line(slide, Inches(4.5), Inches(3.8), Inches(4))
add_text(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.6),
         "Powered by XGBoost Regression", size=18, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.6),
         "Seminar Presentation  |  April 2026", size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ── Slide 2: Table of Contents ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "TABLE OF CONTENTS")
toc = [
    "I.    Introduction",
    "II.   Objective",
    "III.  Dataset Overview & Parameters",
    "IV.   Methodology",
    "V.    Model Results",
    "VI.   Future Scope",
    "VII.  Conclusion",
]
add_bullet_slide(slide, Inches(1.5), Inches(2.2), Inches(10), Inches(5), toc, size=22, color=WHITE)

# ── Slide 3: Introduction ───────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "I. INTRODUCTION")
items = [
    "• Air Quality Index (AQI) is a standardized indicator for reporting daily air quality",
    "• Rising pollution in Indian cities demands real-time, localized monitoring",
    "• Traditional monitoring stations are sparse — Odisha has limited ground-truth coverage",
    "• Machine Learning can bridge the gap using available pollutant & meteorological data",
    "• This project predicts AQI for 360+ locations across 9 districts of Odisha",
    "• Integrates live data from Open-Meteo API for real-time predictions",
]
add_bullet_slide(slide, Inches(1), Inches(2.2), Inches(11), Inches(5), items, size=17)

# ── Slide 4: Objective ──────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "II. OBJECTIVE")
items = [
    "• Build an ML model to accurately predict AQI from pollutant concentrations",
    "• Cover 9 districts: Kalahandi, Dhenkanal, Keonjhar, Khordha, Jajpur, Cuttack, Sundargarh, Ganjam, Sambalpur",
    "• Support 360+ village/town-level locations for granular predictions",
    "• Provide both CLI and Mobile App (React Native/Expo) interfaces",
    "• Integrate live weather & air quality data via Open-Meteo API",
    "• Deploy as a FastAPI backend for scalable real-time inference",
    "• Achieve high prediction accuracy (R² > 0.95) with XGBoost",
]
add_bullet_slide(slide, Inches(1), Inches(2.2), Inches(11), Inches(5), items, size=17)

# ── Slide 5: Dataset Overview ────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "III. DATASET OVERVIEW")
items = [
    "• 20,001 records spanning Jan 2022 – Dec 2025 (4 years)",
    "• 9 districts of Odisha with 360+ unique locations",
    "• 12 columns: Date, Location, District, PM2.5, PM10, NO2, SO2, CO, O3, Temperature, Humidity, AQI, AQI_Category",
    "• No missing values — clean, structured dataset",
    "• AQI Categories: Good, Satisfactory, Moderate, Poor, Very Poor, Severe",
]
add_bullet_slide(slide, Inches(1), Inches(2.2), Inches(11), Inches(4), items, size=17)

# ── Slide 6: Parameters ─────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "III. PARAMETERS (FEATURES)")

# Left column
left_items = [
    "Input Features (8 pollutant/weather):",
    "  • PM2.5 (μg/m³) — Fine particulate matter",
    "  • PM10  (μg/m³) — Coarse particulate matter",
    "  • NO2   (ppb)   — Nitrogen Dioxide",
    "  • SO2   (ppb)   — Sulphur Dioxide",
    "  • CO    (mg/m³) — Carbon Monoxide",
    "  • O3    (ppb)   — Ozone",
    "  • Temperature (°C)",
    "  • Humidity (%)",
]
add_bullet_slide(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(5), left_items, size=16)

right_items = [
    "Engineered Features (5):",
    "  • Month (1–12)",
    "  • Year (2022–2025)",
    "  • DayOfYear (1–365)",
    "  • Location_Encoded (Label)",
    "  • District_Encoded (Label)",
    "",
    "Target Variable:",
    "  • AQI (Air Quality Index, 0–500)",
]
add_bullet_slide(slide, Inches(6.8), Inches(2.2), Inches(5.5), Inches(5), right_items, size=16)



# ── Slide 8: Methodology ────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "IV. METHODOLOGY")
items = [
    "1. Data Collection & Preprocessing",
    "     • Dataset: 20,001 records (2022–2025), 9 districts, 360+ locations",
    "     • Feature engineering: temporal features (Month, Year, DayOfYear)",
    "     • Label Encoding for Location & District categorical variables",
    "",
    "2. Model Selection — XGBoost Regressor",
    "     • Gradient boosting ensemble method — handles non-linear relationships",
    "     • GridSearchCV for hyperparameter tuning (3-fold CV)",
    "     • Parameters tuned: n_estimators, max_depth, learning_rate, subsample, colsample_bytree",
    "",
    "3. Train/Test Split: 80/20 (random_state=42)",
    "",
    "4. Evaluation Metrics: R², MAE, RMSE, 5-fold Cross-Validation",
]
add_bullet_slide(slide, Inches(1), Inches(2.2), Inches(11), Inches(5), items, size=16)

# ── Slide 9: Architecture ───────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "IV. SYSTEM ARCHITECTURE")
items = [
    "Full-Stack Architecture:",
    "",
    "  [Data Layer]     →  odisha_aqi_data.csv (20K records)",
    "  [Training]       →  XGBoost + GridSearchCV + Scikit-learn",
    "  [Model]          →  Serialized via Joblib (.pkl files)",
    "  [Backend API]    →  FastAPI (Python) + Open-Meteo live data",
    "  [Frontend]       →  React Native / Expo mobile app",
    "  [CLI Interface]  →  Interactive Python CLI with Auto/Manual modes",
    "",
    "Live Data Pipeline:",
    "  • Geopy (Nominatim) → Geocodes location to lat/lon",
    "  • Open-Meteo Air Quality API → PM2.5, PM10, CO, NO2, SO2, O3",
    "  • Open-Meteo Weather API → Temperature, Humidity",
]
add_bullet_slide(slide, Inches(1), Inches(2.2), Inches(11), Inches(5), items, size=15)

# ── Slide 10: Model Results ──────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "V. MODEL RESULTS")

# Left column — Regression
reg_items = [
    "Regression (AQI Prediction)",
    "",
    "  ➤  R² Score :  0.9868 (98.68%)",
    "",
    "  ➤  MAE (Mean Absolute Error) :  0.85",
    "",
    "  ➤  RMSE (Root Mean Square Error) :  2.67",
    "",
    "  ➤  5-Fold CV R² :  0.9815 ± 0.0049",
]
add_bullet_slide(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.5), reg_items, size=18)

# Right column — Classification
cls_items = [
    "Classification (AQI Categories)",
    "",
    "  ➤  Accuracy  :  98.83%",
    "",
    "  ➤  Precision :  0.9879",
    "",
    "  ➤  Recall    :  0.9882",
    "",
    "  ➤  F1-Score  :  0.9881",
]
add_bullet_slide(slide, Inches(7), Inches(2.2), Inches(5.5), Inches(4.5), cls_items, size=18)



# ── Slide 14: Future Scope ──────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "VI. FUTURE SCOPE")
items = [
    "• Expand coverage to all 30 districts of Odisha with more granular data",
    "• Integrate satellite imagery (Sentinel-5P) for remote sensing-based predictions",
    "• Implement deep learning models (LSTM, Transformer) for time-series AQI forecasting",
    "• Build a real-time alert system — notify citizens when AQI crosses hazardous levels",
    "• Add health advisory recommendations based on predicted AQI category",
    "• Deploy as a public web dashboard with interactive maps (Mapbox/Leaflet)",
    "• Incorporate traffic density & industrial emission data for better accuracy",
    "• Partner with Odisha State Pollution Control Board for validated ground-truth data",
]
add_bullet_slide(slide, Inches(1), Inches(2.2), Inches(11), Inches(5), items, size=17)

# ── Slide 15: Conclusion ────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "VII. CONCLUSION")
items = [
    "• Successfully built an ML-based AQI prediction system for Odisha",
    "• XGBoost model achieves R² ≈ 0.99 with MAE ≈ 2.5 — highly accurate",
    "• Covers 360+ locations across 9 districts with village-level granularity",
    "• Full-stack deployment: CLI + FastAPI backend + React Native mobile app",
    "• Live data integration via Open-Meteo API enables real-time predictions",
    "• The system can assist policymakers, health agencies, and citizens",
    "  in making informed decisions about air quality and public health",
]
add_bullet_slide(slide, Inches(1), Inches(2.2), Inches(11), Inches(5), items, size=18)

# ── Slide 16: Thank You ─────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(1), Inches(2.5), Inches(11), Inches(1),
         "THANK YOU", size=48, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_accent_line(slide, Inches(5), Inches(3.7), Inches(3))
add_text(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.8),
         "Questions & Discussion", size=24, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.6),
         "AQI Prediction System — Odisha  |  Powered by XGBoost",
         size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Save
prs.save(OUT)
print(f"PPT saved to: {OUT}")
